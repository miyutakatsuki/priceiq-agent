"""Build editable Phase3_Slides.pptx — every text frame is editable in PowerPoint.

Design: minimal, single accent (red #DC2626), Inter font, generous whitespace.
8 slides for ~5 min talk. Run: python3 build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Tokens ───────────────────────────────────────────────────
ACCENT       = RGBColor(0xDC, 0x26, 0x26)
ACCENT_SOFT  = RGBColor(0xFE, 0xE2, 0xE2)
POSITIVE     = RGBColor(0x16, 0xA3, 0x4A)
POSITIVE_SOFT= RGBColor(0xDC, 0xFC, 0xE7)
WARNING      = RGBColor(0xCA, 0x8A, 0x04)
WARNING_SOFT = RGBColor(0xFE, 0xF9, 0xC3)
INFO         = RGBColor(0x25, 0x63, 0xEB)
INFO_SOFT    = RGBColor(0xDB, 0xEA, 0xFE)
VIOLET       = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_SOFT  = RGBColor(0xED, 0xE9, 0xFE)
FG           = RGBColor(0x09, 0x09, 0x0B)
FG2          = RGBColor(0x52, 0x52, 0x5B)
FG3          = RGBColor(0xA1, 0xA1, 0xAA)
FG4          = RGBColor(0xD4, 0xD4, 0xD8)
BORDER       = RGBColor(0xE4, 0xE4, 0xE7)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BG_MUTE      = RGBColor(0xFA, 0xFA, 0xFA)

FONT         = "Inter"
FONT_MONO    = "JetBrains Mono"
SLIDE_W      = Inches(13.333)
SLIDE_H      = Inches(7.5)
M_LEFT       = Inches(0.9)
M_TOP        = Inches(0.85)
M_RIGHT      = Inches(0.9)
CONTENT_W    = SLIDE_W - M_LEFT - M_RIGHT


def make_pres() -> Presentation:
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def blank_slide(p: Presentation):
    layout = p.slide_layouts[6]  # 'Blank'
    s = p.slides.add_slide(layout)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s


def _i(v):
    """Coerce Emu/Inches/int/float to int Emu (PowerPoint XML can't parse float Emu)."""
    return int(v) if v is not None else 0


def add_text(slide, left, top, width, height, text,
             font=FONT, size=18, bold=False, color=FG,
             tracking=0, align=PP_ALIGN.LEFT, line_spacing=1.4):
    """Add a textbox; returns the text_frame for further tuning."""
    tb = slide.shapes.add_textbox(_i(left), _i(top), _i(width), _i(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return tf


def add_eyebrow(slide, top, text):
    """Small uppercase red label at the top of content."""
    add_text(slide, M_LEFT, top, CONTENT_W, Inches(0.3),
             text.upper(), size=11, bold=True, color=ACCENT,
             tracking=0.14, line_spacing=1.0)


def add_h1(slide, top, text, size=48):
    """Big section title."""
    add_text(slide, M_LEFT, top, CONTENT_W, Inches(2.5),
             text, size=size, bold=True, color=FG, line_spacing=1.1)


def add_lead(slide, top, text, width=Inches(11)):
    """Subtitle / lead paragraph in zinc-600."""
    add_text(slide, M_LEFT, top, width, Inches(2),
             text, size=22, color=FG2, line_spacing=1.45)


def add_footer(slide, slide_num, label):
    """Footer rule + section label + slide number (no team names — those live on cover)."""
    rule_y = Inches(6.95)
    text_y = Inches(7.05)
    line = slide.shapes.add_connector(1, M_LEFT, rule_y, SLIDE_W - M_RIGHT, rule_y)
    line.line.color.rgb = BORDER
    line.line.width = Emu(6350)

    add_text(slide, M_LEFT, text_y, Inches(6), Inches(0.28),
             "PriceIQ", size=10, color=FG3, line_spacing=1.0)
    add_text(slide, SLIDE_W - M_RIGHT - Inches(6), text_y, Inches(6), Inches(0.28),
             f"{label}    {slide_num:02d}", size=10, color=FG3,
             align=PP_ALIGN.RIGHT, line_spacing=1.0)


def add_accent_rule(slide, top, width=Inches(0.6)):
    """48px red horizontal rule."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M_LEFT, top, width, Emu(25400))  # 2pt
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT
    rect.line.fill.background()


def add_footnote(slide, top, text, width=None):
    """Small grey footnote near bottom of content area (above 6.95\" footer)."""
    w = width or CONTENT_W
    add_text(slide, M_LEFT, top, w, Inches(0.5),
             text, size=13, color=FG3, line_spacing=1.45)


# ── Slides ───────────────────────────────────────────────────

def slide_cover(p):
    s = blank_slide(p)
    add_text(s, M_LEFT, Inches(2.6), CONTENT_W, Inches(0.4),
             "JHU CAREY  ·  GENERATIVE AI  ·  PHASE 3",
             size=11, bold=True, color=ACCENT, line_spacing=1.0)
    # Title — extra-large
    add_text(s, M_LEFT, Inches(3.05), CONTENT_W, Inches(1.4),
             "PriceIQ", size=80, bold=True, color=FG, line_spacing=1.0)
    add_text(s, M_LEFT, Inches(4.4), Inches(10), Inches(1),
             "A pricing decision agent grounded in 100,000 real e-commerce orders.",
             size=24, color=FG2, line_spacing=1.4)
    add_text(s, M_LEFT, Inches(6.6), Inches(10), Inches(0.4),
             "Kangchun Sun  ·  Tao Cheng  ·  Maoyuan Li      2026",
             size=13, color=FG3, line_spacing=1.0)
    # Accent stripe at top
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(38100))  # 3pt
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT
    rect.line.fill.background()
    return s


def slide_problem(p):
    s = blank_slide(p)
    add_eyebrow(s, M_TOP, "The problem")
    add_h1(s, Inches(1.35),
           "A correct pricing answer needs four kinds of evidence the LLM cannot fabricate.",
           size=36)

    col_w = Inches(5.5)
    col1_left = M_LEFT
    col2_left = M_LEFT + col_w + Inches(0.5)
    top = Inches(4.2)

    add_text(s, col1_left, top, col_w, Inches(0.3),
             "QUANTITATIVE", size=11, bold=True, color=FG3, tracking=0.12)
    add_text(s, col1_left, top + Inches(0.4), col_w, Inches(1.5),
             "Price elasticity β with a confidence interval — measured, not guessed.",
             size=20, color=FG2, line_spacing=1.5)

    add_text(s, col2_left, top, col_w, Inches(0.3),
             "CONTEXTUAL", size=11, bold=True, color=FG3, tracking=0.12)
    add_text(s, col2_left, top + Inches(0.4), col_w, Inches(1.5),
             "Holiday proximity. Weather. Seasonality. Each is a separate signal.",
             size=20, color=FG2, line_spacing=1.5)

    add_footnote(s, Inches(6.45),
                 "A single prompt cannot run OLS. A multi-step agent with typed tools can.")
    add_footer(s, 2, "Problem")
    return s


def _box(slide, left, top, w, h, label, sub=None, fill=WHITE, border=BORDER,
         label_color=FG, sub_color=FG2, label_size=14, sub_size=11, label_bold=True):
    """Draw a labeled rounded rectangle (architecture node)."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _i(left), _i(top), _i(w), _i(h))
    rect.adjustments[0] = 0.10  # corner radius
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = border
    rect.line.width = Emu(6350)
    rect.shadow.inherit = False
    # text
    tf = rect.text_frame
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(label_size)
    r.font.bold = label_bold
    r.font.color.rgb = label_color
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.line_spacing = 1.2
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = FONT
        r2.font.size = Pt(sub_size)
        r2.font.color.rgb = sub_color
    return rect


def _arrow(slide, x1, y1, x2, y2, color=FG3, width_pt=0.75):
    """Thin connector with arrowhead (using straight connector, no head — append shape if needed)."""
    line = slide.shapes.add_connector(1, _i(x1), _i(y1), _i(x2), _i(y2))
    line.line.color.rgb = color
    line.line.width = Emu(int(9525 * width_pt))  # 1pt = 12700 emu, but 0.75pt ≈ 9525
    return line


def slide_architecture(p):
    s = blank_slide(p)
    add_eyebrow(s, M_TOP, "Architecture")
    add_h1(s, Inches(1.35), "Two agents, five tools, one kill switch.", size=36)

    # Layout: two agent boxes on top row, 5 tool boxes below, output box at bottom
    top_y = Inches(2.7)
    agent_w = Inches(3.0)
    agent_h = Inches(1.05)
    gap_x = Inches(0.55)

    # Center horizontally
    total_top_w = agent_w * 2 + gap_x
    top_x = M_LEFT + (CONTENT_W - total_top_w) / 2

    # Planner
    _box(s, top_x, top_y, agent_w, agent_h,
         "Planner", "Sonnet 4.5  ·  XML few-shot",
         fill=ACCENT_SOFT, border=ACCENT, label_color=FG, sub_color=FG2,
         label_size=16, sub_size=12)
    # Executor
    _box(s, top_x + agent_w + gap_x, top_y, agent_w, agent_h,
         "Executor", "Haiku 4.5  ·  tool_use loop",
         fill=INFO_SOFT, border=INFO, label_color=FG, sub_color=FG2,
         label_size=16, sub_size=12)
    # Arrow Planner → Executor
    _arrow(s,
           top_x + agent_w, top_y + agent_h / 2,
           top_x + agent_w + gap_x, top_y + agent_h / 2,
           color=FG2, width_pt=1.25)
    # JSON plan label (above the arrow)
    add_text(s, top_x + agent_w, top_y + agent_h / 2 - Inches(0.42),
             gap_x, Inches(0.3), "JSON plan",
             size=11, color=FG3, align=PP_ALIGN.CENTER, line_spacing=1.0)

    # 5 tools row
    tool_y = Inches(4.4)
    tool_w = Inches(2.05)
    tool_h = Inches(1.0)
    tool_gap = Inches(0.18)
    total_tool_w = tool_w * 5 + tool_gap * 4
    tool_x_start = M_LEFT + (CONTENT_W - total_tool_w) / 2

    tools = [
        ("SQL",        "Olist SQLite",     ACCENT,   ACCENT_SOFT),
        ("OLS",        "Elasticity β",     WARNING,  WARNING_SOFT),
        ("Demand",     "Holidays + season", VIOLET,  VIOLET_SOFT),
        ("Weather",    "OpenWeather",      INFO,     INFO_SOFT),
        ("Simulator",  "3-scenario CI",    POSITIVE, POSITIVE_SOFT),
    ]
    for i, (lbl, sub, col, soft) in enumerate(tools):
        x = tool_x_start + i * (tool_w + tool_gap)
        _box(s, x, tool_y, tool_w, tool_h, lbl, sub,
             fill=WHITE, border=col, label_color=FG, sub_color=FG2,
             label_size=15, sub_size=11)
        # vertical line from tool top to executor bottom
        ex_bottom_y = top_y + agent_h
        ex_x = top_x + agent_w + gap_x + agent_w / 2
        _arrow(s, ex_x, ex_bottom_y, x + tool_w / 2, tool_y, color=FG4, width_pt=0.5)

    # Output box
    out_y = Inches(5.85)
    out_w = Inches(8.5)
    out_h = Inches(0.7)
    out_x = M_LEFT + (CONTENT_W - out_w) / 2
    _box(s, out_x, out_y, out_w, out_h,
         "Recommendation + verbatim causal caveat", None,
         fill=BG_MUTE, border=BORDER, label_color=FG2,
         label_size=14, label_bold=False)

    add_footnote(s, Inches(6.45),
                 "Sonnet plans, Haiku executes.    $0.029 / query — 53% cheaper than all-Sonnet, 20 pts more accurate than all-Haiku.")
    add_footer(s, 3, "Architecture")
    return s


def _bar_chart(slide, left, top, width, height,
               labels, values, baseline, value_colors, baseline_label):
    """Hand-drawn bar chart with rectangles (more design control than XL_CHART)."""
    left, top, width, height = _i(left), _i(top), _i(width), _i(height)
    n = len(values)
    chart_w = width
    chart_h = height
    # leave space for top labels and bottom axis labels
    top_pad = Inches(0.85)
    bot_pad = Inches(0.5)
    plot_h = chart_h - top_pad - bot_pad
    plot_top = top + top_pad
    plot_bottom = plot_top + plot_h
    # baseline line
    max_val = max(max(values), baseline) * 1.10
    # bar layout
    total_bar_w = chart_w * 0.7
    gap = (chart_w - total_bar_w) / (n + 1)
    bar_w = total_bar_w / n
    for i, (lbl, val, col) in enumerate(zip(labels, values, value_colors)):
        bar_h = _i(plot_h * (val / max_val))
        bar_x = _i(left + gap + i * (bar_w + gap))
        bar_y = _i(plot_bottom - bar_h)
        # bar
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, bar_y, _i(bar_w), bar_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = col
        rect.line.fill.background()
        rect.shadow.inherit = False
        # value above bar
        add_text(slide, bar_x, bar_y - Inches(0.55), bar_w, Inches(0.5),
                 f"+{val:.0f}%", size=22, bold=True, color=FG,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
        # label below
        add_text(slide, bar_x, plot_bottom + Inches(0.1), bar_w, Inches(0.5),
                 lbl, size=12, color=FG2,
                 align=PP_ALIGN.CENTER, line_spacing=1.25)
    # baseline dashed line
    base_y = _i(plot_bottom - plot_h * (baseline / max_val))
    base_line = slide.shapes.add_connector(1, _i(left + Inches(0.1)), base_y,
                                            _i(left + chart_w - Inches(0.1)), base_y)
    base_line.line.color.rgb = FG3
    base_line.line.width = Emu(6350)  # ~0.5pt
    base_line.line.dash_style = 7  # dash
    add_text(slide, left + chart_w - Inches(2.5), base_y - Inches(0.32),
             Inches(2.4), Inches(0.3),
             baseline_label, size=10, color=FG3,
             align=PP_ALIGN.RIGHT, line_spacing=1.0)


def slide_demo(p):
    s = blank_slide(p)
    add_eyebrow(s, M_TOP, "Live demo  ·  garden tools −10%")
    add_h1(s, Inches(1.35),
           "5 tools fire in 31 seconds. Three revenue scenarios.", size=28)

    # ── Left column: big number + supporting facts ────────────
    left_w = Inches(4.5)
    add_text(s, M_LEFT, Inches(2.45), left_w, Inches(2.0),
             "+33%", size=110, bold=True, color=ACCENT, line_spacing=1.0)
    add_text(s, M_LEFT, Inches(4.30), left_w, Inches(0.35),
             "CENTRAL SCENARIO", size=11, bold=True, color=FG3,
             tracking=0.12, line_spacing=1.0)
    add_text(s, M_LEFT, Inches(4.65), left_w, Inches(0.5),
             "Revenue lift, 95% CI +21% to +46%",
             size=18, color=FG2, line_spacing=1.4)

    # Drivers — small list
    add_text(s, M_LEFT, Inches(5.30), left_w, Inches(0.3),
             "DRIVERS", size=11, bold=True, color=FG3,
             tracking=0.12, line_spacing=1.0)
    add_text(s, M_LEFT, Inches(5.65), left_w, Inches(1.2),
             "β = −2.83  ·  highly elastic\n"
             "Mother's Day demand × 1.16\n"
             "Rain headwind × 0.94",
             size=14, color=FG2, line_spacing=1.55)

    # ── Right column: real Streamlit screenshot ────────────────
    img_path = "assets/demo_top_main.png"
    import pathlib
    if pathlib.Path(img_path).exists():
        # width 6.4" → height ≈ 4.46". top 2.20" → bottom 6.66" (above footer 6.95)
        s.shapes.add_picture(
            img_path,
            _i(M_LEFT + Inches(5.0)), _i(Inches(2.20)),
            width=_i(Inches(6.4)),
        )
        # Caption above image (no overlap risk with footer)
        add_text(s, M_LEFT + Inches(5.0), Inches(1.85),
                 Inches(6.4), Inches(0.3),
                 "STREAMLIT  ·  CACHED DEMO TAB", size=10, bold=True, color=FG3,
                 tracking=0.12, line_spacing=1.0)

    add_footer(s, 4, "Demo")
    return s


def _tool_chip(slide, left, top, w, h, name, active=True, color=ACCENT):
    """Small chip showing a tool — solid colored if active, dimmed if skipped."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _i(left), _i(top), _i(w), _i(h))
    rect.adjustments[0] = 0.18
    rect.fill.solid()
    if active:
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        text_color = WHITE
    else:
        rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = FG4
        rect.line.width = Emu(6350)
        # dashed border for "skipped"
        rect.line.dash_style = 7
        text_color = FG3
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = name
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = text_color


def _tools_strip(slide, left, top, total_w, active_indices):
    """Render 5 tool chips at the given Y; only `active_indices` are colored."""
    tools = [
        ("SQL",       ACCENT),
        ("OLS",       WARNING),
        ("Demand",    VIOLET),
        ("Weather",   INFO),
        ("Simulate",  POSITIVE),
    ]
    n = 5
    chip_w = Inches(0.95)
    chip_h = Inches(0.55)
    gap = (total_w - chip_w * n) / (n - 1)
    for i, (name, col) in enumerate(tools):
        x = left + i * (chip_w + gap)
        _tool_chip(slide, x, top, chip_w, chip_h, name,
                   active=(i in active_indices), color=col)


def slide_pvc(p):
    s = blank_slide(p)
    add_eyebrow(s, M_TOP, "Reasoning rigor  ·  Prompt Version Control")
    add_h1(s, Inches(1.35), "Same query, two prompts.", size=36)

    col_w = Inches(5.4)
    col1_left = M_LEFT
    col2_left = M_LEFT + col_w + Inches(0.7)
    top_block = Inches(3.0)

    # v1 column header
    add_text(s, col1_left, top_block, col_w, Inches(0.4),
             "V1", size=11, bold=True, color=FG3, tracking=0.12)
    add_text(s, col1_left, top_block + Inches(0.4), col_w, Inches(0.8),
             "3 of 5 tools called",
             size=22, bold=True, color=FG, line_spacing=1.2)
    # v1 tools strip — only SQL, OLS, Simulate active (skipped Demand + Weather)
    _tools_strip(s, col1_left, top_block + Inches(1.4), col_w, active_indices={0, 1, 4})
    add_text(s, col1_left, top_block + Inches(2.2), col_w, Inches(1.2),
             "Plan skips demand and weather. Answer correct in direction but missing context — Shortcut Bias.",
             size=17, color=FG2, line_spacing=1.5)

    # v2 column header
    add_text(s, col2_left, top_block, col_w, Inches(0.4),
             "V2", size=11, bold=True, color=FG3, tracking=0.12)
    add_text(s, col2_left, top_block + Inches(0.4), col_w, Inches(0.8),
             "5 of 5 tools called",
             size=22, bold=True, color=FG, line_spacing=1.2)
    # v2 tools strip — all 5 active
    _tools_strip(s, col2_left, top_block + Inches(1.4), col_w, active_indices={0, 1, 2, 3, 4})
    add_text(s, col2_left, top_block + Inches(2.2), col_w, Inches(1.2),
             "Same query plus 5 worked examples + 11-category list. Plan completes the full pipeline. +45% cost, complete context.",
             size=17, color=FG2, line_spacing=1.5)

    add_footnote(s, Inches(6.45),
                 "In agent design, the prompt is not a hint — it's a contract.\n"
                 "Without examples, the LLM honors the letter of the spec but not the intent.")
    add_footer(s, 5, "PVC story")
    return s


def _progress_ring(slide, cx, cy, outer_r, ring_w, percent, color, track=BG_MUTE):
    """Donut ring approximation: full track circle + colored arc via PIE shapes."""
    # Track (full ring): outer circle + cutout
    track_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          _i(cx - outer_r), _i(cy - outer_r),
                                          _i(outer_r * 2), _i(outer_r * 2))
    track_outer.fill.solid()
    track_outer.fill.fore_color.rgb = track
    track_outer.line.fill.background()
    track_outer.shadow.inherit = False

    # Approximate the colored arc with a pie (works for most percents)
    # MSO_SHAPE.PIE start at 0°/3-o'clock, sweeps counterclockwise — use BLOCK_ARC for ring effect
    if percent >= 100:
        arc = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      _i(cx - outer_r), _i(cy - outer_r),
                                      _i(outer_r * 2), _i(outer_r * 2))
    else:
        # Use PIE shape — the adjustment controls the angle
        arc = slide.shapes.add_shape(MSO_SHAPE.PIE,
                                      _i(cx - outer_r), _i(cy - outer_r),
                                      _i(outer_r * 2), _i(outer_r * 2))
        # PIE adjustments: [start angle, end angle] in 60000ths of a degree
        # Rotate so 0% starts at top (12 o'clock) — needs adjustment 1 = -90, then sweep
        try:
            arc.adjustments[0] = -90  # start at top (degrees)
            arc.adjustments[1] = -90 + (360 * percent / 100)
        except Exception:
            pass
    arc.fill.solid()
    arc.fill.fore_color.rgb = color
    arc.line.fill.background()
    arc.shadow.inherit = False

    # Inner white circle to make it a ring
    inner_r = outer_r - ring_w
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    _i(cx - inner_r), _i(cy - inner_r),
                                    _i(inner_r * 2), _i(inner_r * 2))
    inner.fill.solid()
    inner.fill.fore_color.rgb = WHITE
    inner.line.fill.background()
    inner.shadow.inherit = False


def slide_eval(p):
    s = blank_slide(p)
    add_eyebrow(s, M_TOP, "Evaluation  ·  50 cases  ·  LLM-as-Judge  ·  3-run consistency")
    add_h1(s, Inches(1.35),
           "92% pass rate at $0.029 per query.", size=34)

    # ── Progress ring on the left (92% of pass rate) ──
    ring_cx = M_LEFT + Inches(2.0)
    ring_cy = Inches(4.6)
    outer_r = Inches(1.4)
    ring_w  = Inches(0.30)
    _progress_ring(s, ring_cx, ring_cy, outer_r, ring_w, 92, ACCENT)
    # PASS RATE label above the number (no collision)
    add_text(s, ring_cx - Inches(1.3), ring_cy - Inches(0.50),
             Inches(2.6), Inches(0.30),
             "PASS RATE", size=10, bold=True, color=FG3,
             align=PP_ALIGN.CENTER, tracking=0.12, line_spacing=1.0)
    # 92% number, centered in ring
    add_text(s, ring_cx - Inches(1.3), ring_cy - Inches(0.20),
             Inches(2.6), Inches(0.85),
             "92%", size=44, bold=True, color=FG,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    # Sub-detail just below the ring
    add_text(s, ring_cx - Inches(1.6), Inches(6.20),
             Inches(3.2), Inches(0.35),
             "46 of 50  ·  judge score 4.3 / 5",
             size=12, color=FG2, align=PP_ALIGN.CENTER, line_spacing=1.0)

    # ── Right column: cost & consistency stacked ──
    right_x = M_LEFT + Inches(5.5)
    col_w   = Inches(6.0)   # was 7 → overflowed; CONTENT_W ≈ 11.5 - 5.5 = 6
    # Cost / pass
    add_text(s, right_x, Inches(2.65), col_w, Inches(0.35),
             "COST PER PASS", size=11, bold=True, color=FG3, tracking=0.12, line_spacing=1.0)
    add_text(s, right_x, Inches(3.00), col_w, Inches(1.05),
             "$0.029", size=56, bold=True, color=FG, line_spacing=1.0)
    add_text(s, right_x, Inches(4.05), col_w, Inches(0.4),
             "vs $0.062 all-Sonnet  ·  vs $0.016 all-Haiku (76% pass)",
             size=13, color=FG2, line_spacing=1.4)

    # Category consistency
    add_text(s, right_x, Inches(4.85), col_w, Inches(0.35),
             "CATEGORY CONSISTENCY", size=11, bold=True, color=FG3, tracking=0.12, line_spacing=1.0)
    add_text(s, right_x, Inches(5.20), col_w, Inches(1.05),
             "97%", size=56, bold=True, color=FG, line_spacing=1.0)
    add_text(s, right_x, Inches(6.20), col_w, Inches(0.35),
             "same query → same category across 3 runs",
             size=13, color=FG2, line_spacing=1.4)

    add_footnote(s, Inches(6.45),
                 "Multicollinearity warning: 100% on relevant cases.    "
                 "Causal caveat: included on every successful run.")
    add_footer(s, 6, "Eval & FinOps")
    return s


def slide_failures(p):
    s = blank_slide(p)
    add_eyebrow(s, M_TOP, "Honest failures")
    add_h1(s, Inches(1.35), "Six logged. Four fixed in this phase.", size=40)

    col_w = Inches(5.5)
    col1_left = M_LEFT
    col2_left = M_LEFT + col_w + Inches(0.5)
    top = Inches(3.4)

    # Fixed
    add_text(s, col1_left, top, col_w, Inches(0.4),
             "FIXED", size=11, bold=True, color=FG3, tracking=0.12)
    add_text(s, col1_left, top + Inches(0.5), col_w, Inches(3.0),
             "F-01  Memory threshold confused mid-loop — raised 8K → 30K, added explicit \"finalize NOW\" instruction.\n\n"
             "F-02  Multicollinearity flipped β sign — reduced 7 predictors to 2, added diagnostic + fallback.",
             size=17, color=FG2, line_spacing=1.55)

    # Open
    add_text(s, col2_left, top, col_w, Inches(0.4),
             "OPEN  ·  ACKNOWLEDGED", size=11, bold=True, color=FG3, tracking=0.12)
    add_text(s, col2_left, top + Inches(0.5), col_w, Inches(3.0),
             "F-06  Out-of-scope query \"weather in Tokyo?\" mapped to telefonia instead of refusing. "
             "Tool 1's found:False recovered the run, ~$0.005 wasted.\n\n"
             "Fix in v3 prompt — deferred.",
             size=17, color=FG2, line_spacing=1.55)

    add_footnote(s, Inches(6.45),
                 "Full log: Failure_Log_Phase2.md    Red-team: RT-01..RT-04 in Phase2_Final_Report §5.")
    add_footer(s, 7, "Failures")
    return s


def slide_qa(p):
    s = blank_slide(p)
    add_text(s, M_LEFT, Inches(2.4), CONTENT_W, Inches(0.4),
             "THANK YOU", size=11, bold=True, color=ACCENT, tracking=0.14)
    add_text(s, M_LEFT, Inches(2.85), CONTENT_W, Inches(1.5),
             "Questions?", size=88, bold=True, color=FG, line_spacing=1.0)
    add_accent_rule(s, Inches(4.4))

    # Mono stack — repo / demo / eval
    lines = [
        ("repo", "22 files  ·  9 .py + 10 .md + 2 .json + 1 .ipynb"),
        ("demo", "streamlit run app.py"),
        ("eval", "eval_suite.run_full_eval(n=50, consistency_runs=3)"),
    ]
    top = Inches(4.95)
    line_h = Inches(0.42)
    for i, (key, val) in enumerate(lines):
        y = top + i * line_h
        add_text(s, M_LEFT, y, Inches(1.2), line_h,
                 key, font=FONT_MONO, size=16, bold=True, color=FG, line_spacing=1.0)
        add_text(s, M_LEFT + Inches(1.3), y, Inches(11), line_h,
                 val, font=FONT_MONO, size=16, color=FG2, line_spacing=1.0)

    add_text(s, M_LEFT, Inches(6.7), CONTENT_W, Inches(0.4),
             "Sun  ·  Cheng  ·  Li      JHU Carey  ·  Generative AI  ·  2026",
             size=12, color=FG3)
    return s


# ── Build ────────────────────────────────────────────────────
if __name__ == "__main__":
    p = make_pres()
    slide_cover(p)
    slide_problem(p)
    slide_architecture(p)
    slide_demo(p)
    slide_pvc(p)
    slide_eval(p)
    slide_failures(p)
    slide_qa(p)
    out = "Phase3_Slides.pptx"
    p.save(out)
    print(f"✅ wrote {out}  ({len(p.slides)} slides)")
