"""Layout audit for any .pptx file.

Checks:
  - shape bounds inside slide canvas (no off-canvas overflow)
  - minimum text size ≥ 10pt (warn <10pt; flag <8pt as critical)
  - reports per-slide font-size distribution + char counts

Run:
    python3 audit_pptx.py Phase3_Slides.pptx

Exit code 0 = clean. Non-zero = issues found.
Used as a final-step check in any PPT-modifying workflow — see CLAUDE.md.
"""

import sys
from collections import Counter
from pptx import Presentation


def audit(path: str) -> int:
    p = Presentation(path)
    SW, SH = p.slide_width, p.slide_height
    in_per_emu = 1 / 914400

    def i(v):
        return v * in_per_emu if v else 0.0

    issues = 0
    print(f"Audit: {path}")
    print(f"Canvas: {i(SW):.2f}\" × {i(SH):.2f}\"\n")

    for n, s in enumerate(p.slides, 1):
        bad = []
        sizes = Counter()
        chars = 0
        for sh in s.shapes:
            try:
                x, y = sh.left or 0, sh.top or 0
                w, h = sh.width or 0, sh.height or 0
            except Exception as e:
                bad.append(f"  ❌ unreadable shape '{sh.name}': {e}")
                continue
            if y + h > SH + 1000:
                bad.append(f"  ⚠️  '{sh.name}' bottom {i(y+h):.2f}\" > {i(SH):.2f}\"")
            if x + w > SW + 1000:
                bad.append(f"  ⚠️  '{sh.name}' right  {i(x+w):.2f}\" > {i(SW):.2f}\"")
            if x < -1000 or y < -1000:
                bad.append(f"  ⚠️  '{sh.name}' off-canvas (x={i(x):.2f}\", y={i(y):.2f}\")")
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        sz = run.font.size.pt if run.font.size else None
                        sizes[sz] += len(run.text or "")
                        chars += len(run.text or "")
                        if sz and sz < 8:
                            bad.append(f"  🔴 critical {sz}pt: '{(run.text or '')[:40]}'")
                        elif sz and sz < 10:
                            bad.append(f"  ⚠️  small {sz}pt: '{(run.text or '')[:40]}'")
        sz_str = ", ".join(f"{int(k)}pt:{v}" for k, v in sorted(sizes.items(), reverse=True) if k)
        status = "✅" if not bad else "❌"
        print(f"Slide {n}: {len(s.shapes)} shapes  {chars} chars  {status}")
        if sz_str:
            print(f"  sizes: {sz_str}")
        for line in bad:
            print(line)
            issues += 1
        print()

    print(f"{'=' * 50}")
    print(f"Total issues: {issues}")
    return 0 if issues == 0 else 1


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else "Phase3_Slides.pptx"
    sys.exit(audit(path))
